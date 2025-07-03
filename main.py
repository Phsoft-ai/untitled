# main.py

import base64
import io
import unicodedata
from fastapi import FastAPI, Response
from pydantic import BaseModel, Field
from typing import List, Optional, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

# --- Pydantic 모델 정의 (기존과 동일) ---
class CenterPosition(BaseModel):
    dx: float
    dy: float

class TextItem(BaseModel):
    id: str
    text: str
    center_position: CenterPosition = Field(..., alias='centerPosition')
    font_size_pt: float = Field(..., alias='fontSizePt')
    measured_height_pt: Optional[float] = Field(None, alias='measuredHeightPt')
    color_value: int = Field(..., alias='colorValue')
    font_weight_bold: bool = Field(..., alias='fontWeightBold')
    font_family: Optional[str] = Field(None, alias='fontFamily')

class CanvasData(BaseModel):
    background_image_bytes: Optional[str] = Field(None, alias='backgroundImageBytes')
    canvas_size: Dict[str, float] = Field(..., alias='canvasSize')
    canvas_aspect_ratio: Optional[float] = Field(None, alias='canvasAspectRatio')
    text_items: List[TextItem] = Field(..., alias='textItems')
    excel_data: List[Dict[str, str]] = Field(..., alias='excelData')

app = FastAPI()

def crop_image_to_ratio(img: Image.Image, target_ratio: float) -> Image.Image:
    # 기존과 동일
    img_width, img_height = img.size
    img_ratio = img_width / img_height

    if abs(img_ratio - target_ratio) < 0.01:
        return img

    if img_ratio > target_ratio:
        new_width = int(target_ratio * img_height)
        offset = (img_width - new_width) // 2
        return img.crop((offset, 0, offset + new_width, img_height))
    else:
        new_height = int(img_width / target_ratio)
        offset = (img_height - new_height) // 2
        return img.crop((0, offset, img_width, offset + new_height))

# --- 🚀 FIX: 함수 로직을 명확하고 정확하게 전면 수정 ---
def add_cards_on_slide(slide, chunk_data, text_items_template, canvas_size, cropped_background_stream, target_card_ratio):
    
    # 1. 고정된 A4 1/4 그리드 정의
    page_width_inch = slide.part.package.presentation_part.presentation.slide_width.inches
    page_height_inch = slide.part.package.presentation_part.presentation.slide_height.inches
    
    grid_width_inch = page_width_inch / 2
    grid_height_inch = page_height_inch / 2

    grid_positions = [
        (0, 0),
        (grid_width_inch, 0),
        (0, grid_height_inch),
        (grid_width_inch, grid_height_inch)
    ]

    # 2. 각 그리드 칸(A4 1/4)에 카드 내용 배치
    for i, card_data in enumerate(chunk_data):
        grid_left_inch, grid_top_inch = grid_positions[i]

        # 3. "가상 캔버스" (배경 이미지가 그려질 영역)의 크기와 위치 계산 (BoxFit: Contain 로직)
        grid_ratio = grid_width_inch / grid_height_inch
        
        final_pic_width_inch = grid_width_inch
        final_pic_height_inch = grid_height_inch
        pic_left_offset_inch = 0
        pic_top_offset_inch = 0

        # 받은 카드의 비율(target_card_ratio)에 맞춰 가상 캔버스 크기 조정
        if target_card_ratio > grid_ratio:
            # 카드가 그리드보다 넓으면, 높이를 줄여서 비율을 맞춤 (상하 여백 발생)
            final_pic_height_inch = grid_width_inch / target_card_ratio
            pic_top_offset_inch = (grid_height_inch - final_pic_height_inch) / 2
        elif target_card_ratio < grid_ratio:
            # 카드가 그리드보다 좁으면, 너비를 줄여서 비율을 맞춤 (좌우 여백 발생)
            final_pic_width_inch = grid_height_inch * target_card_ratio
            pic_left_offset_inch = (grid_width_inch - final_pic_width_inch) / 2

        # 4. 배경 이미지 배치 (찌그러짐 없음)
        if cropped_background_stream:
            cropped_background_stream.seek(0)
            slide.shapes.add_picture(
                cropped_background_stream,
                Inches(grid_left_inch + pic_left_offset_inch),
                Inches(grid_top_inch + pic_top_offset_inch),
                width=Inches(final_pic_width_inch),
                height=Inches(final_pic_height_inch)
            )

        # 5. 텍스트 배치
        # Flutter 캔버스의 픽셀 크기
        canvas_width_px = canvas_size['width']
        canvas_height_px = canvas_size['height']
        
        # 픽셀-인치 변환 비율 (가상 캔버스 기준)
        pixels_per_inch_w = canvas_width_px / final_pic_width_inch
        pixels_per_inch_h = canvas_height_px / final_pic_height_inch
        
        for item_template in text_items_template:
            # 텍스트 내용 결정
            if item_template.id == 'title':
                text_content = card_data.get('name', '')
            elif item_template.id == 'subtitle':
                text_content = card_data.get('group', '')
            else:
                text_content = item_template.text

            # Flutter 좌표계를 PPT 좌표계로 변환
            center_pos_px = item_template.center_position
            
            # 텍스트 박스 중심의 절대 좌표 (Flutter 캔버스 좌상단 기준, 픽셀)
            center_x_abs_px = (canvas_width_px / 2) + center_pos_px.dx
            center_y_abs_px = (canvas_height_px / 2) - center_pos_px.dy

            # 텍스트 박스 크기 (인치)
            font_size_pt = item_template.font_size_pt
            measured_height_pt = item_template.measured_height_pt or font_size_pt
            
            box_width_px = canvas_width_px * 0.95 # 너비는 캔버스의 95%
            box_height_px = (measured_height_pt * (96 / 72)) * 1.2 # 높이는 측정된 높이의 1.2배
            
            box_width_inch = box_width_px / pixels_per_inch_w
            box_height_inch = box_height_px / pixels_per_inch_h

            # 텍스트 박스 좌상단 좌표 (가상 캔버스 좌상단 기준, 인치)
            left_rel_px = center_x_abs_px - (box_width_px / 2)
            top_rel_px = center_y_abs_px - (box_height_px / 2)

            left_rel_inch = left_rel_px / pixels_per_inch_w
            top_rel_inch = top_rel_px / pixels_per_inch_h
            
            # 최종 좌표 (그리드 + 가상 캔버스 오프셋 + 텍스트 상대 좌표)
            final_left_inch = grid_left_inch + pic_left_offset_inch + left_rel_inch
            final_top_inch = grid_top_inch + pic_top_offset_inch + top_rel_inch

            # 텍스트 박스 추가 및 스타일링
            txBox = slide.shapes.add_textbox(
                Inches(final_left_inch),
                Inches(final_top_inch),
                Inches(box_width_inch),
                Inches(box_height_inch)
            )
            
            tf = txBox.text_frame
            tf.margin_bottom = Pt(0)
            tf.margin_top = Pt(0)
            tf.margin_left = Pt(0)
            tf.margin_right = Pt(0)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf.clear()

            p = tf.paragraphs[0]
            p.text = text_content
            p.alignment = PP_ALIGN.CENTER

            font = p.font
            default_font = 'Malgun Gothic'

            if item_template.font_family:
                try:
                    font_name = unicodedata.normalize('NFC', item_template.font_family)
                    font.name = font_name
                except Exception:
                    font.name = default_font
            else:
                font.name = default_font

            font.size = Pt(item_template.font_size_pt)
            font.bold = item_template.font_weight_bold
            color_val = item_template.color_value
            font.color.rgb = RGBColor((color_val >> 16) & 0xFF, (color_val >> 8) & 0xFF, color_val & 0xFF)
# --- FIX END ---

@app.post("/generate-ppt")
async def generate_ppt(data: CanvasData):
    prs = Presentation()
    
    prs.slide_width = Inches(8.27)  # A4 가로
    prs.slide_height = Inches(11.69) # A4 세로
    
    chunk_size = 4
    data_chunks = [data.excel_data[i:i + chunk_size] for i in range(0, len(data.excel_data), chunk_size)]

    # --- 🚀 FIX: target_ratio 결정 로직 수정 ---
    # Flutter에서 보낸 canvas_aspect_ratio를 사용하고, 없으면 A4 1/4 비율로 대체
    page_width_inch = prs.slide_width.inches
    page_height_inch = prs.slide_height.inches
    default_card_ratio = (page_width_inch / 2) / (page_height_inch / 2)
    
    target_ratio = default_card_ratio
    if data.canvas_aspect_ratio is not None and data.canvas_aspect_ratio > 0:
        target_ratio = data.canvas_aspect_ratio
    # --- FIX END ---

    cropped_background_stream = None
    if data.background_image_bytes:
        try:
            img_bytes = base64.b64decode(data.background_image_bytes)
            original_image = Image.open(io.BytesIO(img_bytes))
            # 위에서 결정된 target_ratio로 이미지를 자름
            cropped_image = crop_image_to_ratio(original_image, target_ratio)
            cropped_background_stream = io.BytesIO()
            cropped_image.save(cropped_background_stream, format=original_image.format or 'PNG')
            cropped_background_stream.seek(0)
        except Exception as e:
            print(f"Error processing background image: {e}")
            
    for chunk in data_chunks:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # --- 🚀 FIX: 수정된 함수에 target_ratio 전달 ---
        add_cards_on_slide(
            slide=slide,
            chunk_data=chunk,
            text_items_template=data.text_items,
            canvas_size=data.canvas_size,
            cropped_background_stream=cropped_background_stream,
            target_card_ratio=target_ratio,
        )
        # --- FIX END ---

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return Response(
        content=file_stream.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=generated_A4.pptx"}
    )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
